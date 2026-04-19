import streamlit as st
import google.generativeai as genai
import io, requests, re, random
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# 1. CONFIGURATION
genai.configure(api_key=st.secrets["API_KEY"])

# 2. DESIGN POWERPOINT INTERACTIF (Style CFA Chartres)
def appliquer_style_cfa(slide, titre_texte, est_reponse=False):
    couleur_fond = RGBColor(180, 0, 0) if est_reponse else RGBColor(0, 82, 204)
    bandeau = slide.shapes.add_shape(1, 0, 0, Inches(10), Inches(0.7))
    bandeau.fill.solid()
    bandeau.fill.fore_color.rgb = couleur_fond
    bandeau.line.visible = False
    lisere = slide.shapes.add_shape(1, 0, Inches(0.7), Inches(10), Inches(0.05))
    lisere.fill.solid()
    lisere.fill.fore_color.rgb = RGBColor(255, 102, 0)
    lisere.line.visible = False
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.05), Inches(9), Inches(0.6))
    p = txBox.text_frame.paragraphs[0]
    p.text = ("CORRECTION : " if est_reponse else "DÉFI : ") + titre_texte
    p.font.bold, p.font.size, p.font.color.rgb = True, Pt(22), RGBColor(255, 255, 255)

def ajouter_paire_slides(prs, titre, question_txt, reponse_txt, img_prompt=None):
    # --- SLIDE DÉFI (Bleu) ---
    slide_q = prs.slides.add_slide(prs.slide_layouts[6])
    appliquer_style_cfa(slide_q, titre, False)
    has_img = False
    if img_prompt:
        seed = random.randint(1, 99999)
        url = f"https://image.pollinations.ai/prompt/cute_vibrant_cartoon_illustration_of_{img_prompt.replace(' ', '_')}?width=512&height=512&nologo=true&seed={seed}"
        try:
            img_data = requests.get(url, timeout=7).content
            slide_q.shapes.add_picture(io.BytesIO(img_data), Inches(5.6), Inches(1.2), width=Inches(4))
            has_img = True
        except: pass
    width = Inches(4.8) if has_img else Inches(9)
    txBox = slide_q.shapes.add_textbox(Inches(0.5), Inches(1.2), width, Inches(4))
    tf = txBox.text_frame
    tf.word_wrap = True
    for l in question_txt.split('\n'):
        if l.strip():
            p = tf.add_paragraph()
            p.text = "• " + l.replace('*', '').strip()
            p.font.size, p.font.color.rgb = Pt(16), RGBColor(30, 30, 30)
            p.space_after = Pt(10)

    # --- SLIDE RÉPONSE (Rouge) ---
    slide_r = prs.slides.add_slide(prs.slide_layouts[6])
    appliquer_style_cfa(slide_r, titre, True)
    txBox_r = slide_r.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(4))
    tf_r = txBox_r.text_frame
    tf_r.word_wrap = True
    for l in reponse_txt.split('\n'):
        if l.strip():
            p = tf_r.add_paragraph()
            p.text = "✔ " + l.replace('*', '').strip()
            p.font.size, p.font.color.rgb = Pt(16), RGBColor(0, 100, 0)
            p.space_after = Pt(10)

# 3. INTERFACE
st.set_page_config(page_title="Générateur de cours", layout="wide")
st.title("👨‍🏫 Générateur de cours : Édition Studio Chartres")

if 'liste' not in st.session_state:
    st.session_state.liste = ["BP Boucher", "BP Boulanger", "Bac Pro Maintenance Véhicule", "BTS Maintenance Véhicule", "CAP EPC", "BP Coiffure", "AMLHR"]

with st.sidebar:
    st.header("⚙️ Paramètres")
    nouveau = st.text_input("Nouveau diplôme :")
    if st.button("Ajouter") and nouveau:
        st.session_state.liste.append(nouveau)
        st.rerun()

col1, col2 = st.columns(2)
with col1:
    diplome = st.selectbox("Diplôme visé :", st.session_state.liste)
    sujet = st.text_input("Sujet de la leçon :", placeholder="ex: Le système de freinage ABS")
with col2:
    lieu = st.text_input("Lieu du scénario :", value="Chartres")

if st.button("🚀 GÉNÉRER LE PACK COMPLET"):
    if sujet:
        with st.spinner("Conception pédagogique et illustrations en cours..."):
            moteur = [m.name for m in genai.list_models() if 'generateContent' in m.supported_generation_methods][0]
            model = genai.GenerativeModel(moteur)
            
            prompt = f"""Expert pédagogie CFA Chartres. Crée un cours complet pour {diplome} sur {sujet} à {lieu}. 
            Ton ludique, humour, sans salutations. Respecte le référentiel officiel {diplome}.

            TU DOIS GÉNÉRER EXACTEMENT CES 6 SECTIONS DANS CET ORDRE :
            
            ###
            SECTION: Objectifs et Référentiel
            IMAGE: cartoon book and gears
            QUESTION: Quels sont les objectifs du cours et les codes compétences du référentiel {diplome} ciblés ?
            REPONSE: Détaille ici les objectifs et cite les codes officiels (ex: C1.2, S3.1) du programme.

            ###
            SECTION: L'Accroche
            IMAGE: funny scene at {lieu}
            QUESTION: Imagine un scénario humoristique à {lieu} qui pose le problème du jour.
            REPONSE: Explication du problème et de l'enjeu métier.

            ###
            SECTION: Exercice d'Application
            IMAGE: happy professional worker cartoon
            QUESTION: Propose un exercice concret (calcul, diagnostic, geste technique) pour les {diplome}.
            REPONSE: Corrigé détaillé étape par étape.

            ###
            SECTION: Quiz Interactif
            IMAGE: cartoon quiz question mark
            QUESTION: Propose un QCM de 3 questions avec choix A, B, C.
            REPONSE: Les bonnes réponses avec explications.

            ###
            SECTION: Vrai ou Faux
            IMAGE: cartoon checkmark and cross
            QUESTION: Propose 5 affirmations à valider (Vrai ou Faux).
            REPONSE: Correction argumentée pour chaque point.

            ###
            SECTION: Synthèse et Autonomie
            IMAGE: cartoon lightbulb
            QUESTION: Résume l'essentiel et propose une activité de 60 minutes à faire seul.
            REPONSE: Points clés à retenir et consignes de l'activité.
            """
            
            res = model.generate_content(prompt).text
            
            tab1, tab2 = st.tabs(["📝 Document Word (Copier/Coller)", "📥 Télécharger le PowerPoint"])
            
            with tab1:
                st.info("Structure complète pour ton document Word.")
                display_txt = re.sub(r'IMAGE:.*', '', res).replace('SECTION:', '##').replace('QUESTION:', '### Défi Apprenti :').replace('REPONSE:', '### Corrigé Formateur :')
                st.markdown(display_txt)
            
            with tab2:
                prs = Presentation()
                prs.slide_width, prs.slide_height = Inches(10), Inches(5.625)
                duels = res.split('###')
                for d in duels:
                    if "SECTION:" in d:
                        try:
                            titre = re.search(r"SECTION:(.*)", d).group(1).strip()
                            img = re.search(r"IMAGE:(.*)", d).group(1).strip()
                            ques = re.search(r"QUESTION:([\s\S]*?)REPONSE:", d).group(1).strip()
                            rep = re.search(r"REPONSE:([\s\S]*)", d).group(1).strip()
                            ajouter_paire_slides(prs, titre, ques, rep, img)
                        except: pass
                
                buf = io.BytesIO()
                prs.save(buf)
                buf.seek(0)
                st.download_button("📥 Télécharger le PowerPoint Interactif", buf, f"Cours_{sujet}.pptx")
